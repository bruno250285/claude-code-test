import os
import sys
import requests
from dotenv import load_dotenv, set_key
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ENV_FILE = os.path.join(os.path.dirname(__file__), ".env")
load_dotenv(ENV_FILE)

CLIENT_ID = os.getenv("STRAVA_CLIENT_ID")
CLIENT_SECRET = os.getenv("STRAVA_CLIENT_SECRET")
REFRESH_TOKEN = os.getenv("STRAVA_REFRESH_TOKEN")

TOKEN_URL = "https://www.strava.com/oauth/token"
API_BASE = "https://www.strava.com/api/v3"

STRAVA_ORANGE = RGBColor(0xFC, 0x4C, 0x02)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x22, 0x22, 0x22)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)

WMO_CODES = {
    0: "Clear sky", 1: "Mainly clear", 2: "Partly cloudy", 3: "Overcast",
    45: "Fog", 48: "Icy fog",
    51: "Light drizzle", 53: "Drizzle", 55: "Heavy drizzle",
    61: "Light rain", 63: "Rain", 65: "Heavy rain",
    71: "Light snow", 73: "Snow", 75: "Heavy snow",
    80: "Light showers", 81: "Showers", 82: "Heavy showers",
    95: "Thunderstorm", 96: "Thunderstorm w/ hail", 99: "Thunderstorm w/ heavy hail",
}


def abort(msg):
    print(f"Error: {msg}")
    sys.exit(1)


def get_token():
    if not CLIENT_ID or not CLIENT_SECRET or not REFRESH_TOKEN:
        abort("Missing Strava credentials in .env. Run strava_summary.py first to authenticate.")
    resp = requests.post(TOKEN_URL, data={
        "client_id": CLIENT_ID,
        "client_secret": CLIENT_SECRET,
        "grant_type": "refresh_token",
        "refresh_token": REFRESH_TOKEN,
    })
    if not resp.ok:
        abort(f"Token refresh failed: {resp.text}")
    data = resp.json()
    set_key(ENV_FILE, "STRAVA_ACCESS_TOKEN", data["access_token"])
    set_key(ENV_FILE, "STRAVA_REFRESH_TOKEN", data["refresh_token"])
    return data["access_token"]


def get_activities(token, count=5):
    resp = requests.get(
        f"{API_BASE}/athlete/activities",
        headers={"Authorization": f"Bearer {token}"},
        params={"per_page": count},
    )
    if not resp.ok:
        abort(f"Strava API failed: {resp.text}")
    return resp.json()


def get_weather(lat, lng, date_str, hour):
    resp = requests.get(
        "https://archive-api.open-meteo.com/v1/archive",
        params={
            "latitude": lat,
            "longitude": lng,
            "start_date": date_str,
            "end_date": date_str,
            "hourly": "temperature_2m,precipitation,windspeed_10m,weathercode",
            "wind_speed_unit": "kmh",
        },
    )
    if not resp.ok:
        return None
    data = resp.json().get("hourly", {})
    return {
        "temp": data["temperature_2m"][hour],
        "precip": data["precipitation"][hour],
        "wind": data["windspeed_10m"][hour],
        "condition": WMO_CODES.get(data["weathercode"][hour], "Unknown"),
    }


def format_distance(meters):
    return f"{meters / 1000:.2f} km" if meters >= 1000 else f"{int(meters)} m"


def format_duration(seconds):
    h, rem = divmod(int(seconds), 3600)
    m, s = divmod(rem, 60)
    return f"{h}h {m:02d}m {s:02d}s" if h else f"{m}m {s:02d}s"


def add_text(tf, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK_GRAY


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, STRAVA_ORANGE)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    add_text(tf, "My Latest Strava Activities", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.6))
    tf2 = sub_box.text_frame
    add_text(tf2, "Last 5 rides with weather data", 20, color=WHITE, align=PP_ALIGN.CENTER)


def add_activity_slide(prs, activity, index, weather):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)

    # Orange header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = STRAVA_ORANGE
    bar.line.fill.background()

    # Activity number badge
    num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.7), Inches(0.9))
    add_text(num_box.text_frame, f"#{index}", 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Activity name in header
    name_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.25), Inches(8.5), Inches(0.9))
    tf = name_box.text_frame
    tf.word_wrap = True
    add_text(tf, activity.get("name", "Activity"), 28, bold=True, color=WHITE)

    date_local = activity.get("start_date_local", "")
    date = date_local[:10]

    stats = [
        ("Type",        activity.get("sport_type", activity.get("type", "—"))),
        ("Date",        date),
        ("Distance",    format_distance(activity.get("distance", 0))),
        ("Moving Time", format_duration(activity.get("moving_time", 0))),
        ("Elevation",   f"{activity.get('total_elevation_gain', 0):.0f} m gain"),
    ]

    if weather:
        stats.append(("Condition", weather["condition"]))
        stats.append(("Temperature", f"{weather['temp']:.1f} °C"))
        stats.append(("Wind", f"{weather['wind']:.0f} km/h"))
        stats.append(("Precipitation", f"{weather['precip']:.1f} mm"))
    else:
        stats.append(("Weather", "No GPS / unavailable"))

    col_x = [Inches(0.5), Inches(5.2)]
    row_start_y = Inches(1.65)
    row_h = Inches(0.58)

    for i, (label, value) in enumerate(stats):
        col = i // 5
        row = i % 5
        x = col_x[col]
        y = row_start_y + row * row_h

        label_box = slide.shapes.add_textbox(x, y, Inches(1.6), Inches(0.45))
        add_text(label_box.text_frame, label.upper(), 9, bold=True, color=MID_GRAY)

        val_box = slide.shapes.add_textbox(x, y + Inches(0.22), Inches(4.4), Inches(0.35))
        add_text(val_box.text_frame, value, 18, bold=True, color=DARK_GRAY)


def main():
    print("Fetching Strava activities...")
    token = get_token()
    activities = get_activities(token, count=5)

    if not activities:
        abort("No activities found.")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    for i, activity in enumerate(activities, start=1):
        date_local = activity.get("start_date_local", "")
        date = date_local[:10]
        hour = int(date_local[11:13]) if len(date_local) >= 13 else 12
        latlng = activity.get("start_latlng")

        weather = None
        if latlng and len(latlng) == 2:
            print(f"  Fetching weather for activity {i}: {activity.get('name')}...")
            weather = get_weather(latlng[0], latlng[1], date, hour)

        add_activity_slide(prs, activity, i, weather)

    out_path = os.path.join(os.path.dirname(__file__), "strava_activities.pptx")
    prs.save(out_path)
    print(f"\nPresentation saved to: {out_path}")


if __name__ == "__main__":
    main()
